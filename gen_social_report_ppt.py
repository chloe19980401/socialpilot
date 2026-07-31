# -*- coding: utf-8 -*-
"""
生成《社媒运营管理平台 · 7月总结 + 8月规划(SMART)》两页可编辑 PPTX，
样式对齐王力安防月报（金色标题线 + 灰底表头 + 居中表格）。

用法（本机装了 Python 即可）：
    python gen_social_report_ppt.py
会自动安装 python-pptx（若未装），并在脚本同目录生成：
    社媒运营_7月总结_8月规划.pptx
"""
import os, sys, subprocess
try:
    from pptx import Presentation
except ImportError:
    print("未检测到 python-pptx，正在自动安装…")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

FONT = "Microsoft YaHei"                 # Mac 可改 "PingFang SC"
GOLD = RGBColor(0xB8, 0x86, 0x2A)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x8A, 0x8A, 0x8A)
HEADER = RGBColor(0xD9, 0xD9, 0xD9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
REDLBL = RGBColor(0x9C, 0x27, 0x2E)      # 第一列标签用的深红（对齐截图里的红字）

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW = prs.slide_width


def _run(p, text, size, color, bold=False, align=None):
    if align is not None:
        p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color
    r.font.bold = bold; r.font.name = FONT
    return r


def report_slide(title, subtitle, section, headers, rows, label_color=REDLBL):
    s = prs.slides.add_slide(BLANK)
    # 标题
    tb = s.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(9.5), Inches(0.7))
    _run(tb.text_frame.paragraphs[0], title, 26, DARK, True)
    # 金色标题线
    from pptx.enum.shapes import MSO_SHAPE
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.57), Inches(1.02), Inches(3.2), Pt(3))
    line.fill.solid(); line.fill.fore_color.rgb = GOLD; line.line.fill.background(); line.shadow.inherit = False
    # 副标题
    sb = s.shapes.add_textbox(Inches(0.57), Inches(1.12), Inches(9.5), Inches(0.4))
    _run(sb.text_frame.paragraphs[0], subtitle, 12.5, GRAY)
    # 右上角品牌
    bt = s.shapes.add_textbox(Inches(10.4), Inches(0.4), Inches(2.5), Inches(0.4))
    _run(bt.text_frame.paragraphs[0], "王力安防 · WONLY", 12, GOLD, True, PP_ALIGN.RIGHT)
    # 小节标签
    lb = s.shapes.add_textbox(Inches(0.6), Inches(1.75), Inches(6), Inches(0.45))
    _run(lb.text_frame.paragraphs[0], section, 15, DARK, True)

    # 表格
    nrows = len(rows) + 1
    left, top = Inches(0.6), Inches(2.35)
    width, height = Inches(12.1), Inches(0.82) * nrows
    tbl = s.shapes.add_table(nrows, 2, left, top, width, height).table
    tbl.columns[0].width = Inches(3.4)
    tbl.columns[1].width = Inches(8.7)
    tbl.first_row = True
    for i in range(nrows):
        tbl.rows[i].height = Inches(0.82)
    # 表头
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = HEADER
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        _run(c.text_frame.paragraphs[0], h, 15, DARK, True, PP_ALIGN.CENTER)
    # 数据行
    for i, (k, v) in enumerate(rows, start=1):
        c0 = tbl.cell(i, 0)
        c0.fill.solid(); c0.fill.fore_color.rgb = WHITE
        c0.vertical_anchor = MSO_ANCHOR.MIDDLE
        _run(c0.text_frame.paragraphs[0], k, 13.5, label_color, True, PP_ALIGN.CENTER)
        c1 = tbl.cell(i, 1)
        c1.fill.solid(); c1.fill.fore_color.rgb = WHITE
        c1.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = c1.text_frame; tf.word_wrap = True
        _run(tf.paragraphs[0], v, 13, DARK, False, PP_ALIGN.CENTER)
    return s


# ---------- 第 1 页：7月总结 ----------
report_slide(
    "二、社媒运营管理平台",
    "从 0 搭建管理系统到多平台接入 · 王力安防海外社媒 · 2026年7月",
    "总结：",
    ["关键指标", "7月完成"],
    [
        ("管理系统上线", "2026.7 · SocialPilot 平台（内容中心 / 发布排期 / 设计台 / 绩效 / 竞品）"),
        ("平台接入", "YouTube · Instagram · TikTok · Facebook（4 平台 API 数据自动同步）"),
        ("账号粉丝", "YouTube 6 · Instagram 3 · TikTok 1 · Facebook 0"),
        ("内容 & 排期", "已录帖子 10 条 · 内容排期 22 条（已审核 21）"),
        ("自动化", "飞书群日报 + 上传即时提醒 · Meta 永不过期令牌打通"),
        ("当前阶段", "系统就绪期 · 内容起量启动"),
    ],
)

# ---------- 第 2 页：8月规划 SMART ----------
report_slide(
    "二、社媒运营管理平台",
    "阶段目标：从「系统就绪」推进到「内容起量 · 数据增长」",
    "8月规划（SMART）：",
    ["关键指标", "8月需完成"],
    [
        ("发帖量", "4 平台合计 ≥ 40 条（各平台 ≥ 8 条 / 月）"),
        ("月曝光量", "4 平台播放 / 展示合计 约 420 → ≥ 3,000"),
        ("粉丝增长", "四平台合计 10 → ≥ 80"),
        ("按时发布率", "内容排期 ≥ 90% 按时发布 · 逾期 0"),
        ("数据自动化", "IG / FB 帖子级互动量接入（按链接匹配）· 4 平台数据日更"),
        ("重点项目", "KOL 红人库 0 → ≥ 10 家 · 短视频（Reels / TikTok）试产"),
    ],
)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "社媒运营_7月总结_8月规划.pptx")
prs.save(out)
print("已生成：", out)
