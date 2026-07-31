# -*- coding: utf-8 -*-
"""
生成《Meta（IG/FB）Token 长期自动化方案》可编辑 PPTX。

用法（本机装了 Python 即可）：
    python gen_meta_token_ppt.py

脚本会自动安装依赖 python-pptx（若未安装），并在脚本所在目录生成：
    Meta_Token_长期自动化方案.pptx
生成的是原生文本框 / 表格，PowerPoint / WPS 里可直接编辑。
"""

import os
import sys
import subprocess

# ---------- 依赖自检 ----------
try:
    from pptx import Presentation
except ImportError:
    print("未检测到 python-pptx，正在自动安装…")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- 主题 ----------
FONT = "Microsoft YaHei"          # Windows 自带中文字体；Mac 可改 "PingFang SC"
BRAND = RGBColor(0x3B, 0x6E, 0xF6)
DARK = RGBColor(0x1F, 0x29, 0x37)
GRAY = RGBColor(0x64, 0x74, 0x8B)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x22, 0xA5, 0x5A)
RED = RGBColor(0xE1, 0x1D, 0x48)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _set(run, size, color, bold=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT


def textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def band(slide, top, height, color):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, SW, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def header(slide, title, kicker=None):
    band(slide, 0, Inches(1.15), BRAND)
    tf = textbox(slide, Inches(0.6), Inches(0.18), Inches(12), Inches(0.8),
                 MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title; _set(r, 26, WHITE, True)
    if kicker:
        tf2 = textbox(slide, Inches(0.62), Inches(1.25), Inches(12), Inches(0.5))
        r = tf2.paragraphs[0].add_run(); r.text = kicker; _set(r, 13, GRAY)


def bullets(slide, items, top=Inches(1.55), left=Inches(0.7),
            width=Inches(12), size=16, gap=6):
    """items: [(text, level)]  level 0/1"""
    tf = textbox(slide, left, top, width, SH - top - Inches(0.4))
    first = True
    for text, level in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.level = level
        mark = "•  " if level == 0 else "–  "
        r = p.add_run(); r.text = mark + text
        _set(r, size - level * 2, DARK if level == 0 else GRAY, level == 0)
        if level == 1:
            p.space_before = Pt(2)


def table_slide(title, headers, rows, kicker=None, colw=None):
    slide = prs.slides.add_slide(BLANK)
    header(slide, title, kicker)
    nrows, ncols = len(rows) + 1, len(headers)
    left, top = Inches(0.7), Inches(1.7)
    width, height = Inches(12), Inches(0.5) * nrows
    tbl = slide.shapes.add_table(nrows, ncols, left, top, width, height).table
    if colw:
        for i, w in enumerate(colw):
            tbl.columns[i].width = Inches(w)
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = BRAND
        para = c.text_frame.paragraphs[0]; para.alignment = PP_ALIGN.LEFT
        rr = para.add_run(); rr.text = h; _set(rr, 14, WHITE, True)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid(); c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
            para = c.text_frame.paragraphs[0]
            rr = para.add_run(); rr.text = str(val); _set(rr, 12.5, DARK)
    return slide


# ---------- 1. 封面 ----------
s = prs.slides.add_slide(BLANK)
band(s, 0, SH, DARK)
band(s, Inches(2.6), Inches(0.06), BRAND)
tf = textbox(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.6))
r = tf.paragraphs[0].add_run()
r.text = "Meta（Instagram / Facebook）"; _set(r, 34, WHITE, True)
p = tf.add_paragraph(); r = p.add_run()
r.text = "Token 长期自动化方案"; _set(r, 34, BRAND, True)
tf = textbox(s, Inches(0.92), Inches(4.5), Inches(11), Inches(1))
r = tf.paragraphs[0].add_run()
r.text = "让 Graph API 令牌长期有效、无需人工更换，IG/FB 数据持续自动同步"
_set(r, 15, RGBColor(0xC7, 0xD2, 0xE0))
tf = textbox(s, Inches(0.92), Inches(6.4), Inches(11), Inches(0.5))
r = tf.paragraphs[0].add_run()
r.text = "SocialPilot AI · 社媒管理系统"; _set(r, 13, GRAY)

# ---------- 2. 现状诊断 ----------
table_slide(
    "一、现状诊断",
    ["平台", "账号同步", "结论"],
    [
        ["Instagram", "成功", "Token 有效。帖子显示 0 是「手动帖存链接短码、API 存媒体数字 id」对不上，属匹配问题，非 Token。"],
        ["Facebook", "持续失败", "取不到主页数据，多为 Token 过期 / 权限不足 / Page external_id 未填。"],
    ],
    kicker="IG 与 FB 共用同一套 Graph Token，一套方案同时覆盖两者",
    colw=[2.0, 2.2, 7.8],
)

# ---------- 3. 为什么会失效 ----------
s = prs.slides.add_slide(BLANK)
header(s, "二、为什么 Meta Token 会失效")
bullets(s, [
    ("短期用户令牌：1~2 小时过期", 0),
    ("长效用户令牌：60 天过期", 0),
    ("主页（Page）令牌：由长效用户令牌派生，本身长期有效……", 0),
    ("但一旦源用户改密码 / 退出授权 / 60 天链条断掉，即失效", 1),
    ("结论：靠「个人账号手动生成的令牌」必然周期性失效 —— 这正是 FB 现在挂掉的根因", 0),
])

# ---------- 4. 方案 A ----------
s = prs.slides.add_slide(BLANK)
header(s, "三、方案 A（推荐）：系统用户 + 永不过期 Token", "零维护，无需任何续期代码")
bullets(s, [
    ("在 business.facebook.com → 商务设置 → 用户 → 系统用户，新建系统用户（管理员）", 0),
    ("把 Facebook 主页 + Instagram 账号分配给它，权限勾全（管理 + 查看数据）", 0),
    ("生成令牌，勾选权限，过期时间选「永不 Never」", 0),
    ("pages_show_list, pages_read_engagement, read_insights", 1),
    ("instagram_basic, instagram_manage_insights, business_management", 1),
    ("令牌配到后端密钥（只在服务端，绝不进前端/仓库）", 0),
    ("FB_PAGE_TOKEN = 该令牌；IG_ACCESS_TOKEN = 同一个即可；确认 IG_BUSINESS_ID", 1),
    ("确认 accounts 表里 FB/IG 账号的 external_id 填了正确的 Page ID / IG Business ID", 0),
], size=15, gap=5)

# ---------- 5. 方案 B ----------
s = prs.slides.add_slide(BLANK)
header(s, "四、方案 B（备选）：60 天令牌 + 定时自动续期", "复用现有 TikTok 的 tokens 表 + 自动刷新模式")
bullets(s, [
    ("建表 meta_tokens（仅 service_role 可读）：存 access_token + expires_at", 0),
    ("新 Edge Function：meta-token-refresh", 0),
    ("调 Graph 的 fb_exchange_token 接口换新的 60 天令牌，写回表", 1),
    ("pg_cron 每周触发一次（远早于 60 天，避免踩线）", 0),
    ("改 sync-instagram / sync-facebook：Token 从 meta_tokens 表读，回退 env", 0),
    ("需你提供并配为后端密钥：Meta App ID / App Secret + 一个初始长效令牌", 0),
], size=15)

# ---------- 6. A vs B 对比 ----------
table_slide(
    "五、A vs B 对比",
    ["维度", "方案 A：系统用户永不过期", "方案 B：60 天自动续期"],
    [
        ["维护成本", "零，一次配置长期有效", "需搭表+函数+定时任务"],
        ["是否绑个人账号", "否（系统用户）", "是（依赖用户授权链）"],
        ["失效风险", "极低", "低（续期正常即可）"],
        ["需要开发", "不需要", "需要（我可代做）"],
        ["推荐度", "强烈推荐", "备选（无法用 A 时）"],
    ],
    colw=[2.6, 4.9, 4.5],
)

# ---------- 7. 帖子级刷新附带问题 ----------
s = prs.slides.add_slide(BLANK)
header(s, "六、附带问题：IG/FB 帖子级互动量刷新")
bullets(s, [
    ("即使 Token 修好，IG/FB 单条帖子互动量仍刷不到", 0),
    ("手动上传存的 external_id 是链接短码（如 DbaDNt0MWEi）", 1),
    ("Graph API 返回的是媒体数字 id（如 18164672368459388），两者对不上", 1),
    ("解决：把 IG/FB 的刷新改成「按帖子链接 permalink 匹配」（标准化去 query 后比对）", 0),
    ("独立的一小步，可在 Token 通了之后单独做", 1),
])

# ---------- 8. 建议路线 ----------
s = prs.slides.add_slide(BLANK)
header(s, "七、建议路线")
bullets(s, [
    ("① 先走方案 A（系统用户永不过期令牌）—— 最省事，FB 立刻能同步，IG 更稳", 0),
    ("② 若坚持 60 天令牌，再按方案 B 搭自动续期（表 + 函数 + 定时任务）", 0),
    ("③ 最后补 IG/FB 帖子按链接匹配刷新", 0),
], size=17, gap=10)
band(s, Inches(5.5), Inches(1.3), LIGHT)
tf = textbox(s, Inches(0.7), Inches(5.6), Inches(12), Inches(1.1), MSO_ANCHOR.MIDDLE)
r = tf.paragraphs[0].add_run()
r.text = "当前决定：采用方案 A —— 系统用户 + 永不过期 Token"
_set(r, 16, BRAND, True)

# ---------- 保存 ----------
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "Meta_Token_长期自动化方案.pptx")
prs.save(out)
print("已生成：", out)
